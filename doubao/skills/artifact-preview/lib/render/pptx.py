"""PPTX text + LibreOffice→PDF screenshot rendering.

Text via python-pptx; screenshots via the LibreOffice headless converter
that emits a temporary PDF rasterized through :mod:`.pdf`.

LibreOffice missing → empty page list + warning string. Callers can
attach the warning to a manifest without inspecting logs.
"""

from __future__ import annotations

import io
import logging
import os
import subprocess
import tempfile

from .. import _common
from .constants import (
    DEFAULT_LIBREOFFICE_TIMEOUT_SEC,
    DEFAULT_MAX_PAGES_RENDER,
    DEFAULT_PDF_DPI,
)
from .pdf import extract_pdf
from .types import PageImage

logger = logging.getLogger(__name__)


def find_libreoffice_binary(explicit: str | None = None) -> str | None:
    """Return the absolute path to LibreOffice or ``None`` if missing.

    Searches PATH and, because the macOS and Windows installers do not put the
    binary on PATH, the standard install locations for the platform as well.
    ``explicit`` (from ``--soffice``) and ``$ARTIFACT_PREVIEW_SOFFICE`` win.
    """
    return _common.find_office_binary(explicit)


def extract_pptx_text(pptx_bytes: bytes) -> tuple[str, list[str], int]:
    """Extract per-slide text + slide titles from a PPTX byte string.

    Returns
    -------
    (text, slide_titles, slide_count)
        ``text`` is concatenation of ``--- Slide N/M ---`` blocks.
        ``slide_titles`` is the title-placeholder string for each slide
        (``"(untitled)"`` when the layout has no title placeholder).

    Falls back to ``("", [], 0)`` when python-pptx is missing or the
    bytes won't parse.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — pptx text skipped")
        return "", [], 0

    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        logger.warning("[pptx] text open failed: %s", exc)
        return "", [], 0

    text_parts: list[str] = []
    slide_titles: list[str] = []
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        slide_texts: list[str] = []

        # python-pptx exposes title placeholder directly; works for both
        # TITLE and CENTER_TITLE layouts. (The earlier
        # ``placeholder_format.type == TITLE`` check missed CENTER_TITLE
        # because the enum has TWO members for what's effectively the
        # same role.)
        title_text = ""
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        if title_shape is not None and getattr(title_shape, "text", "").strip():
            title_text = title_shape.text.strip().splitlines()[0][:120]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        slide_titles.append(title_text or "(untitled)")
        if slide_texts:
            text_parts.append(
                f"--- Slide {i + 1}/{total_slides} ---\n" + "\n".join(slide_texts)
            )
    return "\n\n".join(text_parts), slide_titles, total_slides


def render_pptx_via_libreoffice(
    pptx_bytes: bytes,
    *,
    dpi: int = DEFAULT_PDF_DPI,
    max_pages: int = DEFAULT_MAX_PAGES_RENDER,
    timeout_sec: float = DEFAULT_LIBREOFFICE_TIMEOUT_SEC,
    soffice_path: str | None = None,
) -> tuple[list[PageImage], str | None]:
    """Convert PPTX → PDF (LibreOffice headless) → PNG pages.

    Returns
    -------
    (pages, error)
        ``pages`` is the rendered slide PNGs (empty on any failure).
        ``error`` is ``None`` on success or a short human-readable
        diagnostic — callers can attach this to a ``warnings`` list
        without having to inspect logs.

    The function never raises: capability gaps (no LibreOffice,
    timeout, bad PPTX, etc.) become an empty page list + a warning
    string, and the caller decides how to surface it.
    """
    lo_bin = find_libreoffice_binary(soffice_path)
    if not lo_bin:
        return [], _common.office_missing_msg()

    with tempfile.TemporaryDirectory(prefix="ar_pptx_") as tmpdir:
        tmp_outdir = tmpdir
        pptx_path = os.path.join(tmp_outdir, "slides.pptx")
        with open(pptx_path, "wb") as fh:
            fh.write(pptx_bytes)
        try:
            result = subprocess.run(
                [
                    lo_bin, "--headless",
                    # as_uri(), not "file://" + path: on Windows the latter
                    # yields file://C:\... where the drive letter is parsed as
                    # the URL host, and LibreOffice then hangs until timeout.
                    "-env:UserInstallation="
                    + _common.path_to_file_uri(os.path.join(tmp_outdir, ".cfg")),
                    "--convert-to", "pdf",
                    "--outdir", tmp_outdir,
                    pptx_path,
                ],
                capture_output=True,
                timeout=timeout_sec,
                check=False,
            )
        except subprocess.TimeoutExpired:
            return [], f"libreoffice timed out after {timeout_sec:.0f}s"
        except Exception as exc:
            return [], f"libreoffice failed: {exc}"

        pdf_path = os.path.join(tmp_outdir, "slides.pdf")
        if not os.path.exists(pdf_path):
            stderr = (result.stderr or b"").decode(errors="replace")[:200]
            return [], f"libreoffice produced no pdf: {stderr or '<no stderr>'}"

        with open(pdf_path, "rb") as fh:
            pdf_bytes = fh.read()

    _text, pages, _total = extract_pdf(
        pdf_bytes,
        dpi=dpi,
        max_pages=max_pages,
        render_images=True,
    )
    return pages, None
